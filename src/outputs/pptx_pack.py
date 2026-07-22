"""
PowerPoint Executive Governance Pack Generator using python-pptx.
Generates an 8-slide executive review pack adhering to PMO presentation standards.
"""
import os
from typing import List, Dict, Any
from src.models.initiative import Initiative
from src.metrics.engine import MetricsEngine

class PowerPointPackGenerator:
    def generate_pack(self, initiatives: List[Initiative], summary_text: str, output_path: str) -> str:
        engine = MetricsEngine()
        m = engine.calculate_portfolio_summary(initiatives)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            blank_slide_layout = prs.slide_layouts[6]

            # Helper function for header + watermark
            def add_slide_header(slide, title_text):
                tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 78, 120)

                wm = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
                wp = wm.text_frame.paragraphs[0]
                wp.text = "SYNTHETIC DEMONSTRATION DATA — NOT FROM A REAL ORGANISATION"
                wp.font.size = Pt(10)
                wp.font.italic = True
                wp.font.color.rgb = RGBColor(127, 127, 127)

            # Slide 1: Title
            s1 = prs.slides.add_slide(blank_slide_layout)
            tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
            tf1 = tb1.text_frame
            p1 = tf1.paragraphs[0]
            p1.text = "Enterprise Executive Portfolio Review Pack"
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(31, 78, 120)
            
            p1_sub = tf1.add_paragraph()
            p1_sub.text = "Monthly Governance & Delivery Progress Update — March 2025"
            p1_sub.font.size = Pt(20)
            p1_sub.font.color.rgb = RGBColor(89, 89, 89)

            wm1 = s1.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5))
            wm1.text_frame.paragraphs[0].text = "SYNTHETIC DEMONSTRATION DATA — NOT FROM A REAL ORGANISATION"

            # Slide 2: Executive Summary
            s2 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s2, "Slide 2: Executive Portfolio Summary")
            tb2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = summary_text
            p2.font.size = Pt(16)

            # Slide 3: Portfolio RAG & Health Overview
            s3 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s3, "Slide 3: Portfolio RAG & Health Overview")
            rows = len(m["by_category"]) + 1
            cols = 5
            table_shape3 = s3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5))
            table3 = table_shape3.table
            headers3 = ["Category", "Initiative Count", "Budget (£M)", "Actual Spend (£M)", "RED Status Count"]
            for col_idx, text in enumerate(headers3):
                table3.cell(0, col_idx).text = text

            for row_idx, (cat, val) in enumerate(m["by_category"].items(), start=1):
                table3.cell(row_idx, 0).text = cat
                table3.cell(row_idx, 1).text = str(val["count"])
                table3.cell(row_idx, 2).text = f"£{val['budget']/1e6:.2f}M"
                table3.cell(row_idx, 3).text = f"£{val['actual_cost']/1e6:.2f}M"
                table3.cell(row_idx, 4).text = str(val["red_count"])

            # Slide 4: Decisions Required
            s4 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s4, "Slide 4: Decisions Required (Ask / Options / Recommendation)")
            tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf4 = tb4.text_frame
            tf4.word_wrap = True
            p4 = tf4.paragraphs[0]
            p4.text = "Decision 1: INIT-001 (API Gateway Acceleration)\n• Ask: Approve additional £250k budget variance.\n• Options: Option A (Extend 3 months); Option B (Inject contractor capacity - Recommended).\n• Recommendation: Option B."
            p4.font.size = Pt(16)

            # Slide 5: Key Risks & Issues
            s5 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s5, "Slide 5: Key Risks & Issues")
            tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf5 = tb5.text_frame
            tf5.word_wrap = True
            tf5.paragraphs[0].text = "• RISK-001: Technical resource contention across Tech & Cyber portfolios (High Impact).\n• ISSUE-001: Data governance approval block on Customer Data Mesh schema (Critical Block)."

            # Slide 6: Dependencies & Blockers
            s6 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s6, "Slide 6: Cross-Portfolio Dependencies")
            tb6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf6 = tb6.text_frame
            tf6.paragraphs[0].text = "• DEP-001: INIT-016 (SOC Automation) requires IAM Zero Trust baseline from INIT-005."

            # Slide 7: Demand & Shaping Pipeline
            s7 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s7, "Slide 7: Demand & Shaping Pipeline Funnel")
            tb7 = s7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf7 = tb7.text_frame
            tf7.paragraphs[0].text = f"Front-Door Funnel Status:\n• Intake Proposals: {m['demand']['intake_count']}\n• Shaping Business Cases: {m['demand']['shaping_count']}\n• Approved Mobilising: {m['demand']['approved_count']}"

            # Slide 8: Next Steps & Governance Cadence
            s8 = prs.slides.add_slide(blank_slide_layout)
            add_slide_header(s8, "Slide 8: Actions & Next Steps")
            tb8 = s8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.8))
            tf8 = tb8.text_frame
            tf8.paragraphs[0].text = "1. Submit Investment Committee decisions for approval by March 25.\n2. Complete monthly data quality exception follow-ups.\n3. Conduct Q2 Portfolio Shaping Review."

            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            prs.save(output_path)
            return output_path
        except ImportError:
            # Fallback if python-pptx is not installed
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path + ".txt", "w", encoding="utf-8") as f:
                f.write(f"POWERPOINT GOVERNANCE PACK (8-SLIDE TEXT SUMMARY)\n\n{summary_text}")
            return output_path + ".txt"
